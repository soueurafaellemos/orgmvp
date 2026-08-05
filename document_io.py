from __future__ import annotations

import io
import mimetypes
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Iterable

import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation


@dataclass
class InputDocument:
    name: str
    data: bytes
    mime_type: str
    original_data: bytes | None = None
    original_mime_type: str | None = None


SUPPORTED_EXTENSIONS = {
    '.pdf',
    '.txt',
    '.md',
    '.json',
    '.html',
    '.htm',
    '.xml',
    '.csv',
    '.tsv',
    '.xls',
    '.xlsx',
    '.docx',
    '.pptx',
    '.eml',
}

LEGACY_EXTENSIONS = {'.doc', '.ppt', '.rtf', '.odt'}
MAX_TEXT_CHARS = 700_000
MAX_ROWS_PER_SHEET = 10_000


def guess_mime(name: str, fallback: str = 'application/octet-stream') -> str:
    guessed, _ = mimetypes.guess_type(name)
    return guessed or fallback


def _decode_text(data: bytes) -> str:
    for encoding in ('utf-8-sig', 'utf-8', 'latin-1'):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode('utf-8', errors='replace')


def _limit_text(text: str, source_name: str) -> str:
    if len(text) <= MAX_TEXT_CHARS:
        return text
    return (
        text[:MAX_TEXT_CHARS]
        + f'\n\n[AVISO: conteúdo de {source_name} truncado após '
          f'{MAX_TEXT_CHARS:,} caracteres.]'
    )


def _converted_text_document(
    *,
    original_name: str,
    original_data: bytes,
    original_mime_type: str,
    text: str,
) -> InputDocument:
    content = _limit_text(text.strip(), original_name)
    return InputDocument(
        name=original_name,
        data=content.encode('utf-8'),
        mime_type='text/plain',
        original_data=original_data,
        original_mime_type=original_mime_type,
    )


def _extract_html(data: bytes) -> str:
    soup = BeautifulSoup(_decode_text(data), 'html.parser')
    return soup.get_text('\n', strip=True)


def _dataframe_to_tsv(df: pd.DataFrame, sheet_name: str) -> str:
    safe = df.copy()
    safe = safe.head(MAX_ROWS_PER_SHEET)
    safe.insert(0, '__linha_planilha__', range(2, len(safe) + 2))
    return (
        f'=== ABA: {sheet_name} ===\n'
        + safe.to_csv(sep='\t', index=False)
    )


def _extract_spreadsheet(data: bytes, name: str) -> str:
    suffix = Path(name).suffix.lower()
    buffer = io.BytesIO(data)
    sections = [
        f'ARQUIVO ORIGINAL: {name}',
        'TIPO: PLANILHA',
        'As colunas estão separadas por tabulação. '
        'A coluna __linha_planilha__ indica a linha aproximada no Excel.',
    ]

    if suffix in {'.csv', '.tsv'}:
        separator = '\t' if suffix == '.tsv' else None
        try:
            df = pd.read_csv(
                buffer,
                sep=separator,
                engine='python',
                dtype=str,
                keep_default_na=False,
            )
        except Exception:
            sections.append(_decode_text(data))
            return '\n\n'.join(sections)
        sections.append(_dataframe_to_tsv(df, 'Planilha'))
        return '\n\n'.join(sections)

    try:
        sheets = pd.read_excel(
            buffer,
            sheet_name=None,
            dtype=str,
            keep_default_na=False,
        )
    except Exception as exc:
        raise ValueError(
            f'Não foi possível ler a planilha {name}: {exc}'
        ) from exc

    for sheet_name, df in sheets.items():
        sections.append(_dataframe_to_tsv(df, str(sheet_name)))

    return '\n\n'.join(sections)


def _extract_docx(data: bytes, name: str) -> str:
    doc = Document(io.BytesIO(data))
    sections = [f'ARQUIVO ORIGINAL: {name}', 'TIPO: DOCUMENTO WORD']

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            sections.append(text)

    for index, table in enumerate(doc.tables, start=1):
        sections.append(f'=== TABELA {index} ===')
        for row in table.rows:
            sections.append('\t'.join(cell.text.strip() for cell in row.cells))

    return '\n'.join(sections)


def _extract_pptx(data: bytes, name: str) -> str:
    presentation = Presentation(io.BytesIO(data))
    sections = [f'ARQUIVO ORIGINAL: {name}', 'TIPO: APRESENTAÇÃO']

    for slide_index, slide in enumerate(presentation.slides, start=1):
        sections.append(f'=== SLIDE {slide_index} ===')
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                sections.append(shape.text.strip())
            if getattr(shape, 'has_table', False):
                for row in shape.table.rows:
                    sections.append(
                        '\t'.join(cell.text.strip() for cell in row.cells)
                    )

    return '\n'.join(sections)


def extract_eml(data: bytes, original_name: str) -> list[InputDocument]:
    message = BytesParser(policy=policy.default).parsebytes(data)
    output: list[InputDocument] = []
    body_parts: list[str] = []

    if message.is_multipart():
        for part in message.walk():
            disposition = part.get_content_disposition()
            content_type = part.get_content_type()

            if disposition == 'attachment':
                filename = part.get_filename() or 'anexo_sem_nome'
                payload = part.get_payload(decode=True) or b''
                suffix = Path(filename).suffix.lower()
                if suffix in SUPPORTED_EXTENSIONS or suffix in LEGACY_EXTENSIONS:
                    output.extend(
                        prepare_documents([(filename, payload, content_type)])
                    )
                continue

            if content_type == 'text/plain':
                try:
                    body_parts.append(part.get_content())
                except Exception:
                    raw = part.get_payload(decode=True) or b''
                    body_parts.append(_decode_text(raw))
    else:
        try:
            body_parts.append(message.get_content())
        except Exception:
            body_parts.append(_decode_text(data))

    headers = [
        f'Assunto: {message.get("subject", "")}',
        f'De: {message.get("from", "")}',
        f'Para: {message.get("to", "")}',
        f'Data: {message.get("date", "")}',
    ]
    body = '\n'.join(headers + ['', *body_parts]).strip()
    if body:
        output.insert(
            0,
            _converted_text_document(
                original_name=original_name,
                original_data=data,
                original_mime_type='message/rfc822',
                text=body,
            ),
        )
    return output


def prepare_documents(
    uploaded: Iterable[tuple[str, bytes, str | None]],
) -> list[InputDocument]:
    docs: list[InputDocument] = []

    for name, data, supplied_mime in uploaded:
        suffix = Path(name).suffix.lower()
        original_mime = supplied_mime or guess_mime(name)

        if suffix in LEGACY_EXTENSIONS:
            raise ValueError(
                f'O formato {suffix} de {name} precisa ser convertido antes '
                'do upload. Salve como .docx, .pptx, .xlsx ou PDF.'
            )

        if suffix not in SUPPORTED_EXTENSIONS:
            raise ValueError(f'Formato não suportado: {name}')

        if suffix == '.eml':
            docs.extend(extract_eml(data, name))
            continue

        if suffix == '.pdf':
            docs.append(
                InputDocument(
                    name=name,
                    data=data,
                    mime_type='application/pdf',
                    original_data=data,
                    original_mime_type=original_mime,
                )
            )
            continue

        if suffix in {'.xlsx', '.xls', '.csv', '.tsv'}:
            text = _extract_spreadsheet(data, name)
        elif suffix == '.docx':
            text = _extract_docx(data, name)
        elif suffix == '.pptx':
            text = _extract_pptx(data, name)
        elif suffix in {'.html', '.htm'}:
            text = _extract_html(data)
        else:
            text = _decode_text(data)

        docs.append(
            _converted_text_document(
                original_name=name,
                original_data=data,
                original_mime_type=original_mime,
                text=text,
            )
        )

    return docs


def split_pdf(
    doc: InputDocument,
    pages_per_batch: int,
    start_page: int = 1,
    end_page: int | None = None,
) -> list[tuple[InputDocument, int, int]]:
    source = fitz.open(stream=doc.data, filetype='pdf')
    total_pages = source.page_count
    start_idx = max(0, start_page - 1)
    end_idx = total_pages if end_page is None else min(total_pages, end_page)

    if start_idx >= end_idx:
        source.close()
        raise ValueError('Intervalo de páginas inválido.')

    batches = []
    for first in range(start_idx, end_idx, pages_per_batch):
        last_exclusive = min(first + pages_per_batch, end_idx)
        part = fitz.open()
        part.insert_pdf(source, from_page=first, to_page=last_exclusive - 1)
        part_bytes = part.tobytes(garbage=4, deflate=True)
        part.close()

        first_human = first + 1
        last_human = last_exclusive
        batches.append(
            (
                InputDocument(
                    name=(
                        f'{Path(doc.name).stem}_'
                        f'p{first_human}-{last_human}.pdf'
                    ),
                    data=part_bytes,
                    mime_type='application/pdf',
                    original_data=doc.original_data or doc.data,
                    original_mime_type=(
                        doc.original_mime_type or 'application/pdf'
                    ),
                ),
                first_human,
                last_human,
            )
        )

    source.close()
    return batches


def render_pdf_page(
    doc: InputDocument,
    page_number: int,
    zoom: float = 1.5,
) -> bytes:
    if doc.mime_type != 'application/pdf':
        raise ValueError('A fonte selecionada não é um PDF.')

    pdf = fitz.open(stream=doc.data, filetype='pdf')
    try:
        page_index = int(page_number) - 1
        if page_index < 0 or page_index >= pdf.page_count:
            raise ValueError(f'Página {page_number} fora do intervalo do PDF.')
        page = pdf.load_page(page_index)
        pixmap = page.get_pixmap(
            matrix=fitz.Matrix(zoom, zoom),
            alpha=False,
        )
        return pixmap.tobytes('png')
    finally:
        pdf.close()
