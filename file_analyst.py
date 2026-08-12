from __future__ import annotations

"""NAVE Intelligence Core — File Analyst v1.

O File Analyst trata o ARQUIVO como uma unidade de inteligência completa.
Ele não substitui a análise de projeto; ele cria a camada de evidência que torna
as análises de projeto, cliente, fornecedor, venue e portfólio auditáveis.

Princípios:
- evidência antes da inferência;
- proveniência estável em página/slide/parágrafo/linha/célula/imagem;
- nenhuma regra específica de cliente/projeto;
- falha de IA nunca apaga a evidência determinística;
- orçado/proposto nunca é tratado como realizado;
- o pipeline legado pode continuar rodando em paralelo (dual-write).
"""

import hashlib
import io
import json
import mimetypes
import os
import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

from pydantic import BaseModel, Field


FILE_ANALYST_VERSION = "file-analyst-v1"
FILE_ANALYST_SCHEMA_VERSION = "1"
FILE_ANALYST_PROMPT_VERSION = "2026-08-11.v1"
DEFAULT_MODEL = "gemini-3.5-flash-lite"

ENTITY_TYPES = {
    "project", "client", "brand", "supplier", "venue", "venue_space",
    "product", "concept", "strategy", "activation", "solution", "deliverable",
    "requirement", "audience", "platform", "technology", "journey_stage",
    "gift", "presskit", "communication_asset", "financial_line_item",
    "cost_category", "person", "partner", "kpi", "report", "campaign",
    "event", "location", "risk", "learning",
}

PREDICATES = {
    "budget_max", "expected_attendees", "event_date", "start_date", "end_date",
    "capacity", "area_sqm", "price", "proposed_total", "actual_total", "quantity",
    "duration", "deadline", "satisfaction", "sentiment", "commercial_result",
    "execution_result", "approval_status", "geographic_coverage", "lead_time_days",
    "cost_per_attendee", "preferred_format", "required_platform_behavior",
    "decision_reason", "constraint_status",
}

RELATION_TYPES = {
    "requirement_of", "responds_to", "fulfills", "partially_fulfills",
    "does_not_fulfill", "conflicts_with", "materializes", "derived_from_concept",
    "uses_venue", "uses_supplier", "uses_product", "uses_technology",
    "targets_audience", "belongs_to_journey_stage", "creates_content_for_platform",
    "depends_on", "part_of", "delivered_as", "costed_by", "quoted_by",
    "paid_by_client", "optional_cost_of", "budget_constraint_of", "validated_by",
    "challenged_by", "approved_by", "criticized_by", "rejected_by", "modified_to",
    "replaced_by", "executed_as", "similar_to", "variant_of", "inspired_by",
    "repeated_for_client", "historically_associated_with", "performed_with",
    "project_for_client", "belongs_to_brand",
}

ROLE_ALIASES = {
    "briefing_original": "briefing_original",
    "proposal_presentation": "proposal_presentation",
    "final_presentation": "final_presentation",
    "detailed_costs": "cost_sheet",
    "preliminary_budget": "budget_sheet",
    "feedback_approval": "feedback",
    "post_event_report": "post_event_report",
    "supplier_reference": "supplier_reference",
    "complementary_document": "complementary_document",
}


class EvidenceUnit(BaseModel):
    ref: str
    unit_type: str
    ordinal: int | None = None
    locator: dict[str, Any] = Field(default_factory=dict)
    content_text: str | None = None
    content_json: dict[str, Any] = Field(default_factory=dict)
    extraction_method: str = "parser"
    extraction_confidence: float = 1.0
    language: str | None = None


class EntityCandidate(BaseModel):
    key: str
    entity_type: str
    canonical_name: str
    aliases: list[str] = Field(default_factory=list)
    entity_kind: str = "project_instance"
    confidence: float = 0.75
    evidence_refs: list[str] = Field(default_factory=list)
    attributes: dict[str, Any] = Field(default_factory=dict)


class ClaimCandidate(BaseModel):
    subject_key: str
    predicate: str
    value_type: str
    object_key: str | None = None
    value_text: str | None = None
    value_numeric: float | None = None
    value_boolean: bool | None = None
    value_date: str | None = None
    value_timestamp: str | None = None
    value_json: Any | None = None
    unit: str | None = None
    currency: str | None = None
    claim_kind: str = "fact"
    confidence: float = 0.75
    authority_score: float | None = None
    evidence_refs: list[str] = Field(default_factory=list)


class RelationCandidate(BaseModel):
    source_key: str
    relation_type: str
    target_key: str
    relation_kind: str = "fact"
    confidence: float = 0.75
    authority_score: float | None = None
    evidence_refs: list[str] = Field(default_factory=list)
    attributes: dict[str, Any] = Field(default_factory=dict)


class FileSemanticChunk(BaseModel):
    title: str | None = None
    language: str | None = None
    source_role: str | None = None
    source_role_confidence: float = 0.0
    summary: str | None = None
    entities: list[EntityCandidate] = Field(default_factory=list)
    claims: list[ClaimCandidate] = Field(default_factory=list)
    relations: list[RelationCandidate] = Field(default_factory=list)
    unknowns: list[str] = Field(default_factory=list)
    contradictions: list[str] = Field(default_factory=list)
    transcription: str | None = None


class FileAnalysisResult(BaseModel):
    file_name: str
    mime_type: str
    sha256: str
    source_role: str
    source_role_confidence: float
    source_role_reasons: list[str] = Field(default_factory=list)
    title: str | None = None
    language: str | None = None
    summary: str | None = None
    evidence_units: list[EvidenceUnit] = Field(default_factory=list)
    entities: list[EntityCandidate] = Field(default_factory=list)
    claims: list[ClaimCandidate] = Field(default_factory=list)
    relations: list[RelationCandidate] = Field(default_factory=list)
    unknowns: list[str] = Field(default_factory=list)
    contradictions: list[str] = Field(default_factory=list)
    warnings: list[str] = Field(default_factory=list)
    semantic_analysis_ran: bool = False
    metadata: dict[str, Any] = Field(default_factory=dict)
    analyzer_version: str = FILE_ANALYST_VERSION


@dataclass(frozen=True)
class SemanticJob:
    label: str
    docs: list[Any]
    evidence_refs: list[str]
    source_page_offset: int = 0


def _normalize(value: Any) -> str:
    text = unicodedata.normalize("NFKD", str(value or ""))
    text = "".join(ch for ch in text if not unicodedata.combining(ch))
    text = re.sub(r"[^a-z0-9]+", " ", text.casefold())
    return re.sub(r"\s+", " ", text).strip()


def _stable_key(entity_type: str, name: str) -> str:
    base = _normalize(name).replace(" ", "_")[:70] or entity_type
    return f"{entity_type}:{base}"


def _sha(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _json_safe(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, Mapping):
        return {str(k): _json_safe(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [_json_safe(v) for v in value]
    return str(value)


def _clip(value: Any, limit: int) -> str:
    text = str(value or "").strip()
    return text if len(text) <= limit else text[: limit - 1] + "…"


def _guess_mime(name: str, supplied: str | None = None) -> str:
    return str(supplied or mimetypes.guess_type(name)[0] or "application/octet-stream")


def _resolve_ai_settings() -> tuple[str | None, str]:
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    model = os.getenv("GEMINI_MODEL") or DEFAULT_MODEL
    try:
        import streamlit as st  # type: ignore

        api_key = str(
            st.secrets.get("GEMINI_API_KEY")
            or st.secrets.get("GOOGLE_API_KEY")
            or api_key
            or ""
        ).strip() or None
        model = str(st.secrets.get("GEMINI_MODEL") or model).strip() or model
    except Exception:
        pass
    return api_key, model


def _extract_pdf_units(name: str, data: bytes) -> list[EvidenceUnit]:
    import fitz  # lazy: não exige PDF para outros testes

    units: list[EvidenceUnit] = []
    pdf = fitz.open(stream=data, filetype="pdf")
    try:
        for idx in range(pdf.page_count):
            page = pdf.load_page(idx)
            text = re.sub(r"\s+", " ", page.get_text("text") or "").strip()
            image_count = len(page.get_images(full=True))
            ordinal = idx + 1
            units.append(EvidenceUnit(
                ref=f"page:{ordinal}",
                unit_type="page",
                ordinal=ordinal,
                locator={"page": ordinal},
                content_text=_clip(text, 18000) or None,
                content_json={
                    "page_width": round(float(page.rect.width), 2),
                    "page_height": round(float(page.rect.height), 2),
                    "image_count": image_count,
                    "text_char_count": len(text),
                },
                extraction_method="pymupdf_text",
                extraction_confidence=0.98 if text else 0.75,
            ))
    finally:
        pdf.close()
    return units


def _extract_docx_units(name: str, data: bytes) -> list[EvidenceUnit]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    units: list[EvidenceUnit] = []
    ordinal = 0
    for p_index, paragraph in enumerate(doc.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        ordinal += 1
        units.append(EvidenceUnit(
            ref=f"paragraph:{ordinal}",
            unit_type="paragraph",
            ordinal=ordinal,
            locator={"paragraph_index": p_index},
            content_text=_clip(text, 12000),
            content_json={"style": str(getattr(paragraph.style, "name", "") or "")},
            extraction_method="python-docx",
            extraction_confidence=0.99,
        ))
    for t_index, table in enumerate(doc.tables, start=1):
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if not any(any(cell for cell in row) for row in rows):
            continue
        ordinal += 1
        units.append(EvidenceUnit(
            ref=f"table:{t_index}",
            unit_type="table",
            ordinal=ordinal,
            locator={"table": t_index},
            content_text="\n".join("\t".join(row) for row in rows)[:30000] or None,
            content_json={"rows": rows[:120]},
            extraction_method="python-docx",
            extraction_confidence=0.99,
        ))
    return units


def _extract_pptx_units(name: str, data: bytes) -> list[EvidenceUnit]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    units: list[EvidenceUnit] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        image_count = 0
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:  # PICTURE
                image_count += 1
            shape_text = str(getattr(shape, "text", "") or "").strip()
            if shape_text:
                texts.append(shape_text)
            if getattr(shape, "has_table", False):
                table_rows: list[list[str]] = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.strip() for cell in row.cells])
                tables.append(table_rows)
        text = "\n".join(texts).strip()
        units.append(EvidenceUnit(
            ref=f"slide:{idx}",
            unit_type="slide",
            ordinal=idx,
            locator={"slide": idx},
            content_text=_clip(text, 18000) or None,
            content_json={"shape_count": len(slide.shapes), "image_count": image_count, "tables": tables[:8]},
            extraction_method="python-pptx",
            extraction_confidence=0.98 if text else 0.75,
        ))
    return units


def _extract_xlsx_units(name: str, data: bytes) -> list[EvidenceUnit]:
    from openpyxl import load_workbook

    keep_vba = Path(name).suffix.lower() == ".xlsm"
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True, keep_vba=keep_vba)
    units: list[EvidenceUnit] = []
    global_ordinal = 0
    try:
        for ws in wb.worksheets:
            global_ordinal += 1
            units.append(EvidenceUnit(
                ref=f"sheet:{ws.title}",
                unit_type="sheet",
                ordinal=global_ordinal,
                locator={"sheet": ws.title},
                content_json={"max_row": ws.max_row, "max_column": ws.max_column},
                extraction_method="openpyxl",
                extraction_confidence=0.99,
            ))
            blank_streak = 0
            for row_num, values in enumerate(ws.iter_rows(values_only=True), start=1):
                values = list(values[:80])
                if not any(v not in (None, "") and str(v).strip() for v in values):
                    blank_streak += 1
                    if blank_streak > 80 and row_num > 200:
                        break
                    continue
                blank_streak = 0
                cells = {
                    str(index + 1): _json_safe(value)
                    for index, value in enumerate(values)
                    if value not in (None, "") and str(value).strip()
                }
                text = " | ".join(f"C{col}={value}" for col, value in cells.items())
                global_ordinal += 1
                units.append(EvidenceUnit(
                    ref=f"row:{ws.title}:{row_num}",
                    unit_type="row",
                    ordinal=global_ordinal,
                    locator={"sheet": ws.title, "row": row_num},
                    content_text=_clip(text, 12000) or None,
                    content_json={"cells": cells},
                    extraction_method="openpyxl_values",
                    extraction_confidence=0.99,
                ))
    finally:
        wb.close()
    return units


def _extract_text_units(name: str, data: bytes) -> list[EvidenceUnit]:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = data.decode("utf-8", errors="replace")
    chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n+", text) if chunk.strip()]
    if not chunks:
        chunks = [text.strip()] if text.strip() else []
    return [
        EvidenceUnit(
            ref=f"paragraph:{idx}",
            unit_type="paragraph",
            ordinal=idx,
            locator={"paragraph": idx},
            content_text=_clip(chunk, 18000),
            extraction_method="text_decode",
            extraction_confidence=0.98,
        )
        for idx, chunk in enumerate(chunks, start=1)
    ]


def extract_evidence_units(file_name: str, data: bytes, mime_type: str | None = None) -> list[EvidenceUnit]:
    suffix = Path(file_name).suffix.lower()
    mime = _guess_mime(file_name, mime_type)
    if suffix == ".pdf" or mime == "application/pdf":
        return _extract_pdf_units(file_name, data)
    if suffix == ".docx":
        return _extract_docx_units(file_name, data)
    if suffix == ".pptx":
        return _extract_pptx_units(file_name, data)
    if suffix in {".xlsx", ".xlsm"}:
        return _extract_xlsx_units(file_name, data)
    if suffix in {".csv", ".tsv", ".txt", ".md", ".json", ".html", ".htm", ".xml", ".eml"} or mime.startswith("text/"):
        return _extract_text_units(file_name, data)
    if mime.startswith("image/") or suffix in {".jpg", ".jpeg", ".png", ".webp"}:
        return [EvidenceUnit(
            ref="image:1",
            unit_type="image",
            ordinal=1,
            locator={"image": 1},
            content_json={"file_name": file_name, "mime_type": mime, "file_size_bytes": len(data)},
            extraction_method="binary_image",
            extraction_confidence=1.0,
        )]
    return [EvidenceUnit(
        ref="other:1",
        unit_type="other",
        ordinal=1,
        locator={"file": file_name},
        content_json={"mime_type": mime, "file_size_bytes": len(data)},
        extraction_method="binary_metadata",
        extraction_confidence=1.0,
    )]


def _evidence_text(units: Sequence[EvidenceUnit], refs: Sequence[str] | None = None, *, max_chars: int = 90000) -> str:
    wanted = set(refs or [])
    parts: list[str] = []
    total = 0
    for unit in units:
        if wanted and unit.ref not in wanted:
            continue
        text = unit.content_text or ""
        if not text:
            continue
        block = f"[[EVIDENCE:{unit.ref}]]\n{text}\n"
        if total + len(block) > max_chars:
            break
        parts.append(block)
        total += len(block)
    return "\n".join(parts)


def _infer_role(file_name: str, units: Sequence[EvidenceUnit], declared_role: str | None) -> tuple[str, float, list[str]]:
    excerpt = _evidence_text(units, max_chars=120000)
    inferred_role = None
    inferred_conf = 0.0
    reasons: list[str] = []
    try:
        from project_batch_ingestion import classify_document

        inferred_role, inferred_conf, reasons = classify_document(file_name, excerpt)
    except Exception:
        pass

    declared = str(declared_role or "").strip()
    if declared:
        # Papel específico já revisado no lote tem precedência. Só substituímos
        # automaticamente papéis historicamente genéricos/propensos a erro quando
        # a identidade estrutural do arquivo é forte. Isso evita uma planilha de
        # custos virar fornecedor apenas porque possui coluna Vendor/Subcontractor.
        weak_declared = declared in {"complementary_document", "supplier_reference"}
        if inferred_role and inferred_role != declared and (
            (weak_declared and inferred_conf >= 0.92)
            or (not weak_declared and inferred_conf >= 0.97)
        ):
            return inferred_role, float(inferred_conf), list(reasons)
        confidence = max(float(inferred_conf or 0.0), 0.82 if not weak_declared else 0.72)
        if inferred_role and inferred_role != declared:
            reasons = [*reasons, f"papel do lote preservado: {declared}; inferência concorrente: {inferred_role} ({inferred_conf:.2f})"]
        return declared, confidence, reasons or ["papel já classificado no lote"]
    if inferred_role and inferred_conf >= 0.90:
        return inferred_role, float(inferred_conf), list(reasons)
    if inferred_role:
        return inferred_role, float(inferred_conf), list(reasons)

    normalized_name = _normalize(Path(file_name).stem)
    if "feedback" in normalized_name or "retorno" in normalized_name:
        return "feedback_approval", 0.82, ["nome do arquivo sinaliza feedback"]
    return "complementary_document", 0.45, ["papel não resolvido com confiança suficiente"]


def _clean_entity(candidate: EntityCandidate) -> EntityCandidate | None:
    entity_type = str(candidate.entity_type or "").strip()
    name = re.sub(r"\s+", " ", str(candidate.canonical_name or "")).strip()
    if entity_type not in ENTITY_TYPES or not name:
        return None
    candidate.entity_type = entity_type
    candidate.canonical_name = name[:320]
    candidate.key = str(candidate.key or "").strip() or _stable_key(entity_type, name)
    candidate.aliases = list(dict.fromkeys(
        re.sub(r"\s+", " ", str(v)).strip()[:320]
        for v in candidate.aliases
        if str(v).strip()
    ))[:12]
    candidate.entity_kind = candidate.entity_kind if candidate.entity_kind in {
        "canonical", "project_instance", "scoped_profile", "ephemeral"
    } else "project_instance"
    candidate.confidence = max(0.0, min(1.0, float(candidate.confidence or 0.0)))
    candidate.evidence_refs = list(dict.fromkeys(str(v) for v in candidate.evidence_refs if str(v).strip()))[:24]
    return candidate


def _clean_claim(candidate: ClaimCandidate) -> ClaimCandidate | None:
    if candidate.predicate not in PREDICATES:
        return None
    candidate.subject_key = str(candidate.subject_key or "").strip() or "project"
    candidate.value_type = candidate.value_type if candidate.value_type in {
        "entity", "text", "numeric", "boolean", "date", "timestamp", "json"
    } else "text"
    if candidate.value_type == "entity" and not candidate.object_key:
        return None
    if candidate.value_type == "numeric" and candidate.value_numeric is None:
        return None
    if candidate.value_type == "text" and not str(candidate.value_text or "").strip():
        return None
    if candidate.value_type == "boolean" and candidate.value_boolean is None:
        return None
    if candidate.value_type == "date" and not candidate.value_date:
        return None
    if candidate.value_type == "timestamp" and not candidate.value_timestamp:
        return None
    if candidate.value_type == "json" and candidate.value_json is None:
        return None
    candidate.claim_kind = candidate.claim_kind if candidate.claim_kind in {
        "fact", "inference", "decision", "recommendation"
    } else "fact"
    candidate.confidence = max(0.0, min(1.0, float(candidate.confidence or 0.0)))
    if candidate.authority_score is not None:
        candidate.authority_score = max(0.0, min(1.0, float(candidate.authority_score)))
    candidate.evidence_refs = list(dict.fromkeys(str(v) for v in candidate.evidence_refs if str(v).strip()))[:24]
    if candidate.currency:
        candidate.currency = str(candidate.currency).upper()[:3]
    return candidate


def _clean_relation(candidate: RelationCandidate) -> RelationCandidate | None:
    if candidate.relation_type not in RELATION_TYPES:
        return None
    if not candidate.source_key or not candidate.target_key or candidate.source_key == candidate.target_key:
        return None
    candidate.relation_kind = candidate.relation_kind if candidate.relation_kind in {
        "fact", "inference", "decision", "recommendation"
    } else "fact"
    candidate.confidence = max(0.0, min(1.0, float(candidate.confidence or 0.0)))
    if candidate.authority_score is not None:
        candidate.authority_score = max(0.0, min(1.0, float(candidate.authority_score)))
    candidate.evidence_refs = list(dict.fromkeys(str(v) for v in candidate.evidence_refs if str(v).strip()))[:24]
    return candidate


FILE_ANALYST_PROMPT = r"""
Você é o File Analyst da NAVE by VOE — plataforma proprietária de inteligência de
pré-produção para live marketing.

Analise SOMENTE a fonte fornecida neste lote. Não use conhecimento externo e não faça
cruzamento com outros arquivos que você não recebeu. O objetivo é transformar o arquivo
em conhecimento auditável e reutilizável.

REGRAS ABSOLUTAS
1. Evidência antes de inferência. Toda entidade, claim ou relação deve citar evidence_refs
   exatamente no formato [[EVIDENCE:...]] fornecido no texto de apoio.
2. Não invente fornecedor, local, data, valor, execução, aprovação ou resultado.
3. Proposta/orçamento não é gasto real. Só use actual_total/execution_result se a própria
   fonte comprovar execução/realizado.
4. Um arquivo pode conter várias coisas: estratégia, conceito, ativação, requisito, venue,
   plataforma, press kit etc. Extraia entidades distintas, sem jogar tudo em um item genérico.
5. Preserve granularidade: conceito ≠ ativação ≠ jornada ≠ press kit ≠ comunicação.
6. Se houver contradição INTERNA no próprio arquivo, registre em contradictions. Se faltar
   evidência, registre em unknowns.
7. source_role deve ser uma função documental, não um assunto citado incidentalmente.

ENTIDADES PERMITIDAS
project, client, brand, supplier, venue, venue_space, product, concept, strategy,
activation, solution, deliverable, requirement, audience, platform, technology,
journey_stage, gift, presskit, communication_asset, financial_line_item, cost_category,
person, partner, kpi, report, campaign, event, location, risk, learning.

PREDICADOS PERMITIDOS
budget_max, expected_attendees, event_date, start_date, end_date, capacity, area_sqm,
price, proposed_total, actual_total, quantity, duration, deadline, satisfaction, sentiment,
commercial_result, execution_result, approval_status, geographic_coverage, lead_time_days,
cost_per_attendee, preferred_format, required_platform_behavior, decision_reason,
constraint_status.

RELAÇÕES PERMITIDAS
requirement_of, responds_to, fulfills, partially_fulfills, does_not_fulfill,
conflicts_with, materializes, derived_from_concept, uses_venue, uses_supplier,
uses_product, uses_technology, targets_audience, belongs_to_journey_stage,
creates_content_for_platform, depends_on, part_of, delivered_as, costed_by, quoted_by,
paid_by_client, optional_cost_of, budget_constraint_of, validated_by, challenged_by,
approved_by, criticized_by, rejected_by, modified_to, replaced_by, executed_as,
similar_to, variant_of, inspired_by, repeated_for_client, historically_associated_with,
performed_with, project_for_client, belongs_to_brand.

CHAVES
- Use "project" para o projeto atual quando a fonte afirmar algo sobre o projeto como um todo.
- Para outras entidades, gere key curta, estável e sem acentos, como "cinemateca" ou
  "youtube_requirement". A mesma entidade dentro deste lote deve manter a mesma key.
- Entidades globais claras (YouTube, Instagram, TikTok, marcas, venues, fornecedores)
  podem usar entity_kind="canonical". Conceitos/ativações/requisitos específicos do projeto
  devem usar entity_kind="project_instance".

ROLE possíveis
briefing_original, proposal_presentation, final_presentation, detailed_costs,
preliminary_budget, feedback_approval, post_event_report, supplier_reference,
complementary_document.

Retorne SOMENTE JSON válido no schema solicitado.
""".strip()


def _semantic_call(*, docs: list[Any], prompt: str, api_key: str, model: str) -> FileSemanticChunk:
    # Usa a superfície já instalada na NAVE, mantendo import lazy para testes locais.
    from gemini_extractor import _structured_call, get_client

    client = get_client(api_key)
    return _structured_call(
        client,
        model=model,
        prompt=prompt,
        docs=docs,
        schema=FileSemanticChunk,
        context="File Analyst",
    )


def _text_doc(name: str, text: str) -> Any:
    from document_io import InputDocument
    return InputDocument(name=name, data=text.encode("utf-8"), mime_type="text/plain")


def _pdf_jobs(file_name: str, data: bytes, units: Sequence[EvidenceUnit], *, pages_per_job: int = 28) -> list[SemanticJob]:
    from document_io import InputDocument, split_pdf

    source = InputDocument(name=file_name, data=data, mime_type="application/pdf", original_data=data, original_mime_type="application/pdf")
    jobs: list[SemanticJob] = []
    for part, first, last in split_pdf(source, pages_per_batch=pages_per_job):
        refs = [f"page:{page}" for page in range(first, last + 1)]
        text = _evidence_text(units, refs, max_chars=75000)
        guide = (
            f"Este PDF parcial corresponde às páginas ORIGINAIS {first}–{last}.\n"
            "Use SOMENTE os evidence_refs abaixo para citar evidência; os números são das páginas originais.\n\n"
            + text
        )
        jobs.append(SemanticJob(
            label=f"pages_{first}_{last}",
            docs=[_text_doc(f"{Path(file_name).stem}_evidence_p{first}-{last}.txt", guide), part],
            evidence_refs=refs,
            source_page_offset=first - 1,
        ))
    return jobs


def _text_jobs(file_name: str, units: Sequence[EvidenceUnit], *, max_chars: int = 85000) -> list[SemanticJob]:
    jobs: list[SemanticJob] = []
    current_refs: list[str] = []
    current_parts: list[str] = []
    current_len = 0
    for unit in units:
        if not unit.content_text:
            continue
        block = f"[[EVIDENCE:{unit.ref}]]\n{unit.content_text}\n"
        if current_parts and current_len + len(block) > max_chars:
            text = "\n".join(current_parts)
            jobs.append(SemanticJob(
                label=f"text_{len(jobs)+1}",
                docs=[_text_doc(f"{Path(file_name).stem}_evidence_{len(jobs)+1}.txt", text)],
                evidence_refs=list(current_refs),
            ))
            current_refs, current_parts, current_len = [], [], 0
        current_refs.append(unit.ref)
        current_parts.append(block)
        current_len += len(block)
    if current_parts:
        jobs.append(SemanticJob(
            label=f"text_{len(jobs)+1}",
            docs=[_text_doc(f"{Path(file_name).stem}_evidence_{len(jobs)+1}.txt", "\n".join(current_parts))],
            evidence_refs=list(current_refs),
        ))
    return jobs


def _image_jobs(file_name: str, data: bytes, mime_type: str, units: Sequence[EvidenceUnit]) -> list[SemanticJob]:
    from document_io import InputDocument
    guide = _text_doc(
        f"{Path(file_name).stem}_evidence_guide.txt",
        "Use o evidence_ref [[EVIDENCE:image:1]] para qualquer afirmação diretamente suportada pela imagem.",
    )
    image = InputDocument(name=file_name, data=data, mime_type=mime_type, original_data=data, original_mime_type=mime_type)
    return [SemanticJob(label="image_1", docs=[guide, image], evidence_refs=["image:1"])]


def _semantic_jobs(file_name: str, data: bytes, mime_type: str, units: Sequence[EvidenceUnit]) -> list[SemanticJob]:
    suffix = Path(file_name).suffix.lower()
    if suffix == ".pdf" or mime_type == "application/pdf":
        return _pdf_jobs(file_name, data, units)
    if mime_type.startswith("image/") or suffix in {".jpg", ".jpeg", ".png", ".webp"}:
        return _image_jobs(file_name, data, mime_type, units)
    return _text_jobs(file_name, units)


def _merge_semantic_chunks(chunks: Sequence[FileSemanticChunk]) -> tuple[
    list[EntityCandidate], list[ClaimCandidate], list[RelationCandidate], list[str], list[str], str | None, str | None, str | None
]:
    entities_by_key: dict[str, EntityCandidate] = {}
    claims_by_sig: dict[str, ClaimCandidate] = {}
    relations_by_sig: dict[str, RelationCandidate] = {}
    unknowns: list[str] = []
    contradictions: list[str] = []
    titles: list[str] = []
    languages: list[str] = []
    summaries: list[str] = []

    for chunk in chunks:
        if chunk.title:
            titles.append(chunk.title)
        if chunk.language:
            languages.append(chunk.language)
        if chunk.summary:
            summaries.append(chunk.summary)
        unknowns.extend(chunk.unknowns)
        contradictions.extend(chunk.contradictions)

        for raw in chunk.entities:
            entity = _clean_entity(raw)
            if entity is None:
                continue
            key = entity.key
            if key in entities_by_key:
                current = entities_by_key[key]
                current.evidence_refs = list(dict.fromkeys([*current.evidence_refs, *entity.evidence_refs]))[:40]
                current.aliases = list(dict.fromkeys([*current.aliases, *entity.aliases]))[:20]
                current.confidence = max(current.confidence, entity.confidence)
                current.attributes.update(entity.attributes)
            else:
                entities_by_key[key] = entity

        for raw in chunk.claims:
            claim = _clean_claim(raw)
            if claim is None:
                continue
            value_sig = json.dumps({
                "object": claim.object_key,
                "text": claim.value_text,
                "numeric": claim.value_numeric,
                "boolean": claim.value_boolean,
                "date": claim.value_date,
                "timestamp": claim.value_timestamp,
                "json": _json_safe(claim.value_json),
                "unit": claim.unit,
                "currency": claim.currency,
            }, ensure_ascii=False, sort_keys=True)
            sig = f"{claim.subject_key}|{claim.predicate}|{claim.value_type}|{value_sig}"
            if sig in claims_by_sig:
                current = claims_by_sig[sig]
                current.evidence_refs = list(dict.fromkeys([*current.evidence_refs, *claim.evidence_refs]))[:40]
                current.confidence = max(current.confidence, claim.confidence)
            else:
                claims_by_sig[sig] = claim

        for raw in chunk.relations:
            relation = _clean_relation(raw)
            if relation is None:
                continue
            sig = f"{relation.source_key}|{relation.relation_type}|{relation.target_key}"
            if sig in relations_by_sig:
                current = relations_by_sig[sig]
                current.evidence_refs = list(dict.fromkeys([*current.evidence_refs, *relation.evidence_refs]))[:40]
                current.confidence = max(current.confidence, relation.confidence)
                current.attributes.update(relation.attributes)
            else:
                relations_by_sig[sig] = relation

    title = max(titles, key=len) if titles else None
    language = max(set(languages), key=languages.count) if languages else None
    summary = " ".join(dict.fromkeys(summaries))[:5000] if summaries else None
    return (
        list(entities_by_key.values()),
        list(claims_by_sig.values()),
        list(relations_by_sig.values()),
        list(dict.fromkeys(v for v in unknowns if str(v).strip()))[:80],
        list(dict.fromkeys(v for v in contradictions if str(v).strip()))[:80],
        title,
        language,
        summary,
    )



def _parse_br_money(raw: str) -> float | None:
    text = re.sub(r"[^\d,.-]", "", str(raw or "")).strip()
    if not text:
        return None
    if "," in text and "." in text:
        if text.rfind(",") > text.rfind("."):
            text = text.replace(".", "").replace(",", ".")
        else:
            text = text.replace(",", "")
    elif "," in text:
        text = text.replace(".", "").replace(",", ".")
    try:
        return float(text)
    except Exception:
        return None


def _add_briefing_deterministic_semantics(role: str, units: Sequence[EvidenceUnit], result: FileAnalysisResult) -> None:
    if role != "briefing_original":
        return
    platform_names = {
        "youtube": "YouTube", "instagram": "Instagram", "tiktok": "TikTok",
        "kwai": "Kwai", "linkedin": "LinkedIn", "facebook": "Facebook",
        "twitch": "Twitch",
    }
    budget_candidates: list[tuple[float, str]] = []
    attendee_candidates: list[tuple[float, str]] = []
    text_units = [unit for unit in units if unit.content_text]

    # Claims numéricos de alta confiança.
    for unit in text_units:
        text = str(unit.content_text or "")
        norm = _normalize(text)
        if any(token in norm for token in ("budget", "orcamento", "verba", "teto orcamentario", "budget cap")):
            for match in re.finditer(r"(?i)(?:R\$|BRL|\$)\s*([0-9][0-9. ]*(?:,[0-9]{1,2})?)", text):
                value = _parse_br_money(match.group(1))
                if value is not None and value >= 1000:
                    budget_candidates.append((value, unit.ref))
        for match in re.finditer(r"(?i)\b([0-9]{2,5})\s*(?:convidados|participantes|guests?|attendees?)\b", text):
            try:
                value = float(match.group(1))
            except Exception:
                continue
            if 10 <= value <= 100000:
                attendee_candidates.append((value, unit.ref))

    # Direcionais de plataforma são tratados como SEÇÕES, não como proximidade
    # lexical global. Isso evita atribuir um "vídeo horizontal/vertical" genérico
    # do final do briefing à ativação YouTube só porque a palavra YouTube reaparece.
    section_starts: list[tuple[int, str, str]] = []
    for idx, unit in enumerate(text_units):
        norm = _normalize(unit.content_text)
        for token, display in platform_names.items():
            if token in norm and any(anchor in norm for anchor in ("ativacao", "activation", "diretriz", "guideline")):
                section_starts.append((idx, token, display))
                break
    for pos, (start_idx, token, display) in enumerate(section_starts):
        end_idx = section_starts[pos + 1][0] if pos + 1 < len(section_starts) else min(len(text_units), start_idx + 20)
        section = text_units[start_idx:end_idx]
        section_norm = _normalize(" ".join(str(row.content_text or "") for row in section))
        behavior = None
        terms: tuple[str, ...] = ()
        if any(term in section_norm for term in ("horizontal", "16 9", "landscape")):
            behavior = "horizontal"
            terms = ("horizontal", "16 9", "landscape")
        elif any(term in section_norm for term in ("vertical", "verticais", "9 16", "portrait")):
            behavior = "vertical"
            terms = ("vertical", "verticais", "9 16", "portrait")
        elif any(term in section_norm for term in ("lifestyle", "aspiracional", "aspirational")):
            behavior = "lifestyle / aspirational content"
            terms = ("lifestyle", "aspiracional", "aspirational")
        if not behavior:
            continue
        behavior_ref = section[0].ref
        for row in section:
            row_norm = _normalize(row.content_text)
            if any(term in row_norm for term in terms):
                behavior_ref = row.ref
                break
        refs = list(dict.fromkeys([section[0].ref, behavior_ref]))
        platform_key = token
        req_key = f"{token}_requirement"
        result.entities.extend([
            EntityCandidate(key=platform_key, entity_type="platform", canonical_name=display, entity_kind="canonical", confidence=0.99, evidence_refs=refs),
            EntityCandidate(key=req_key, entity_type="requirement", canonical_name=f"{display} — formato/comportamento exigido", entity_kind="project_instance", confidence=0.94, evidence_refs=refs),
        ])
        result.claims.append(ClaimCandidate(
            subject_key=req_key, predicate="required_platform_behavior", value_type="text",
            value_text=behavior, claim_kind="fact", confidence=0.94, authority_score=0.95, evidence_refs=refs
        ))
        result.relations.append(RelationCandidate(
            source_key=req_key, relation_type="requirement_of", target_key="project",
            relation_kind="fact", confidence=0.96, authority_score=0.95, evidence_refs=refs
        ))

    if budget_candidates:
        value, ref = max(budget_candidates, key=lambda item: item[0])
        result.claims.append(ClaimCandidate(
            subject_key="project", predicate="budget_max", value_type="numeric", value_numeric=value,
            currency="BRL", claim_kind="fact", confidence=0.94, authority_score=0.95, evidence_refs=[ref]
        ))
    if attendee_candidates:
        value, ref = max(attendee_candidates, key=lambda item: item[0])
        result.claims.append(ClaimCandidate(
            subject_key="project", predicate="expected_attendees", value_type="numeric", value_numeric=value,
            unit="people", claim_kind="fact", confidence=0.88, authority_score=0.90, evidence_refs=[ref]
        ))

def _add_financial_semantics(file_name: str, data: bytes, role: str, units: Sequence[EvidenceUnit], result: FileAnalysisResult) -> None:
    if role not in {"detailed_costs", "preliminary_budget"}:
        return
    try:
        from memory_cost_parser import parse_cost_workbook

        parsed = parse_cost_workbook(file_name, data)
    except Exception as exc:
        result.warnings.append(f"Parser financeiro do File Analyst não conseguiu estruturar a planilha: {exc}")
        return

    proposed_total = parsed.client_total
    if proposed_total is not None:
        # Evidência preferida: linha com maior probabilidade de TOTAL no mesmo sheet.
        evidence_ref = None
        total_candidates = [
            unit for unit in units
            if unit.unit_type == "row"
            and str(unit.locator.get("sheet") or "") == parsed.sheet_name
            and "total" in _normalize(unit.content_text)
        ]
        if total_candidates:
            evidence_ref = total_candidates[-1].ref
        result.claims.append(ClaimCandidate(
            subject_key="project",
            predicate="proposed_total",
            value_type="numeric",
            value_numeric=float(proposed_total),
            currency=parsed.currency or "BRL",
            claim_kind="fact",
            confidence=0.99,
            authority_score=0.95,
            evidence_refs=[evidence_ref] if evidence_ref else [],
        ))

    for item in parsed.items:
        key = f"financial_line:{_normalize(parsed.sheet_name).replace(' ', '_')}:{int(item.source_row)}"
        name = str(item.item_name or f"Linha {item.source_row}").strip()
        ref = f"row:{parsed.sheet_name}:{int(item.source_row)}"
        result.entities.append(EntityCandidate(
            key=key,
            entity_type="financial_line_item",
            canonical_name=name,
            entity_kind="project_instance",
            confidence=0.99,
            evidence_refs=[ref],
            attributes={
                "source_sheet": parsed.sheet_name,
                "source_row": int(item.source_row),
                "category": item.category,
                "vendor": item.raw_data.get("vendor") if isinstance(item.raw_data, Mapping) else None,
                "description": item.description,
                "billing_type": item.billing_type,
                "item_status": item.item_status,
                "estimate_type": item.estimate_type,
                "base_value": item.base_value,
                "fees_value": item.fees_value,
                "charges_value": item.charges_value,
                "client_total": item.client_total,
            },
        ))
        if item.client_total is not None:
            result.claims.append(ClaimCandidate(
                subject_key=key,
                predicate="price",
                value_type="numeric",
                value_numeric=float(item.client_total),
                currency=parsed.currency or "BRL",
                claim_kind="fact",
                confidence=0.99,
                authority_score=0.95,
                evidence_refs=[ref],
            ))
        if item.quantity is not None:
            result.claims.append(ClaimCandidate(
                subject_key=key,
                predicate="quantity",
                value_type="numeric",
                value_numeric=float(item.quantity),
                claim_kind="fact",
                confidence=0.99,
                authority_score=0.95,
                evidence_refs=[ref],
            ))
        if item.category:
            category_key = _stable_key("cost_category", item.category)
            result.entities.append(EntityCandidate(
                key=category_key,
                entity_type="cost_category",
                canonical_name=item.category,
                entity_kind="project_instance",
                confidence=0.99,
                evidence_refs=[ref],
            ))


def _add_feedback_semantics(file_name: str, data: bytes, mime_type: str, role: str, result: FileAnalysisResult, *, api_key: str | None, model: str) -> bool:
    if role != "feedback_approval":
        return False
    if not api_key:
        return False
    try:
        from project_analyst import analyze_feedback_bytes

        analysis = analyze_feedback_bytes(
            file_name=file_name,
            mime_type=mime_type,
            file_bytes=data,
            api_key=api_key,
            model=model,
        )
    except Exception as exc:
        result.warnings.append(f"Leitura especializada de feedback no File Analyst falhou: {exc}")
        return False

    result.metadata.update({
        "feedback_source_type": analysis.source_type,
        "feedback_process_stage": analysis.process_stage,
        "feedback_process_type": analysis.process_type,
        "feedback_proposal_result": analysis.proposal_result,
        "feedback_execution_result": analysis.execution_result,
        "feedback_confidence_level": analysis.confidence_level,
        "feedback_decision_summary": analysis.decision_summary,
        "feedback_result_reasons": list(analysis.result_reasons),
        "feedback_claims": [claim.model_dump() for claim in analysis.claims],
    })
    if analysis.transcription:
        result.evidence_units.append(EvidenceUnit(
            ref="transcript:1",
            unit_type="transcript_segment",
            ordinal=1,
            locator={"derived_from": "image:1" if mime_type.startswith("image/") else "source"},
            content_text=_clip(analysis.transcription, 30000),
            content_json={"source_type": analysis.source_type, "process_stage": analysis.process_stage},
            extraction_method="gemini_multimodal_transcription",
            extraction_confidence=0.95 if analysis.confidence_level == "client_confirmed" else 0.82,
        ))
    if analysis.commercial_result not in {"in_evaluation", "not_informed", "not_applicable"}:
        result.claims.append(ClaimCandidate(
            subject_key="project",
            predicate="commercial_result",
            value_type="text",
            value_text=analysis.commercial_result,
            claim_kind="decision",
            confidence=0.98 if analysis.confidence_level == "client_confirmed" else 0.82,
            authority_score=1.0 if analysis.source_type in {"client", "procurement", "marketing", "branding"} else 0.85,
            evidence_refs=["transcript:1"] if analysis.transcription else ["image:1"],
        ))
    for index, claim in enumerate(analysis.claims, start=1):
        entity_name = claim.related_entities[0] if claim.related_entities else claim.title
        if not entity_name:
            continue
        # O tipo é conservador e baseado no tema. Entity resolution posterior pode corrigir.
        entity_type = {
            "creative_concept": "concept",
            "strategy": "strategy",
            "activation": "activation",
            "scenography": "solution",
            "gift": "gift",
            "journey": "journey_stage",
            "technology": "technology",
        }.get(claim.theme, "solution")
        if claim.result_reason == "venue" or "venue" in _normalize(claim.title) or "local" in _normalize(claim.title):
            entity_type = "venue"
        key = _stable_key(entity_type, entity_name)
        result.entities.append(EntityCandidate(
            key=key,
            entity_type=entity_type,
            canonical_name=entity_name,
            entity_kind="canonical" if entity_type == "venue" else "project_instance",
            confidence=0.84,
            evidence_refs=["transcript:1"] if analysis.transcription else ["image:1"],
        ))
        result.claims.append(ClaimCandidate(
            subject_key=key,
            predicate="sentiment",
            value_type="text",
            value_text=claim.sentiment,
            claim_kind="decision",
            confidence=0.95,
            authority_score=1.0 if analysis.source_type in {"client", "procurement", "marketing", "branding"} else 0.85,
            evidence_refs=["transcript:1"] if analysis.transcription else ["image:1"],
        ))
        if claim.item_outcome_status not in {"unassessed", "unknown"}:
            result.claims.append(ClaimCandidate(
                subject_key=key,
                predicate="approval_status",
                value_type="text",
                value_text=claim.item_outcome_status,
                claim_kind="decision",
                confidence=0.95,
                authority_score=1.0 if analysis.source_type in {"client", "procurement", "marketing", "branding"} else 0.85,
                evidence_refs=["transcript:1"] if analysis.transcription else ["image:1"],
            ))
        if claim.interpretation or claim.result_reason:
            reason = claim.result_reason or claim.interpretation
            result.claims.append(ClaimCandidate(
                subject_key=key,
                predicate="decision_reason",
                value_type="text",
                value_text=_clip(reason, 2000),
                claim_kind="decision",
                confidence=0.90,
                authority_score=1.0 if analysis.source_type in {"client", "procurement", "marketing", "branding"} else 0.85,
                evidence_refs=["transcript:1"] if analysis.transcription else ["image:1"],
            ))
    return True


def analyze_file(
    *,
    file_name: str,
    data: bytes,
    mime_type: str | None = None,
    declared_role: str | None = None,
    enable_semantic: bool = True,
    api_key: str | None = None,
    model: str | None = None,
) -> FileAnalysisResult:
    """Analisa um único arquivo, sem persistir no banco.

    A função é deliberadamente pura o suficiente para ser usada pelo IQ Bench.
    """
    mime = _guess_mime(file_name, mime_type)
    units = extract_evidence_units(file_name, data, mime)
    role, role_confidence, role_reasons = _infer_role(file_name, units, declared_role)
    result = FileAnalysisResult(
        file_name=file_name,
        mime_type=mime,
        sha256=_sha(data),
        source_role=role,
        source_role_confidence=role_confidence,
        source_role_reasons=role_reasons,
        title=Path(file_name).stem,
        evidence_units=list(units),
    )

    _add_briefing_deterministic_semantics(role, units, result)
    _add_financial_semantics(file_name, data, role, units, result)

    resolved_api_key, resolved_model = _resolve_ai_settings()
    api_key = api_key if api_key is not None else resolved_api_key
    model = model or resolved_model

    # Feedback possui um extrator dedicado mais granular que o File Analyst genérico.
    if enable_semantic and _add_feedback_semantics(file_name, data, mime, role, result, api_key=api_key, model=model):
        result.semantic_analysis_ran = True
    elif enable_semantic and api_key:
        chunks: list[FileSemanticChunk] = []
        try:
            jobs = _semantic_jobs(file_name, data, mime, units)
            for job in jobs:
                prompt = (
                    f"{FILE_ANALYST_PROMPT}\n\n"
                    f"Arquivo: {file_name}\nMIME: {mime}\n"
                    f"Papel previamente resolvido: {role} (confiança {role_confidence:.2f}). "
                    "Você pode discordar somente se a função documental estiver claramente evidenciada.\n\n"
                    f"Lote: {job.label}"
                )
                chunks.append(_semantic_call(docs=job.docs, prompt=prompt, api_key=api_key, model=model))
            if chunks:
                entities, claims, relations, unknowns, contradictions, title, language, summary = _merge_semantic_chunks(chunks)
                result.entities.extend(entities)
                result.claims.extend(claims)
                result.relations.extend(relations)
                result.unknowns.extend(unknowns)
                result.contradictions.extend(contradictions)
                result.title = title or result.title
                result.language = language
                result.summary = summary
                result.semantic_analysis_ran = True
        except Exception as exc:
            result.warnings.append(f"Análise semântica do File Analyst não pôde ser concluída; evidência determinística preservada: {exc}")
    elif enable_semantic and not api_key:
        result.warnings.append("File Analyst semântico não executado porque GEMINI_API_KEY não está configurada; evidências determinísticas foram preservadas.")

    # Revalida e deduplica inclusive os candidatos determinísticos adicionados depois do merge.
    entities_by_key: dict[str, EntityCandidate] = {}
    for raw in result.entities:
        entity = _clean_entity(raw)
        if entity is None:
            continue
        key = entity.key
        if key in entities_by_key:
            current = entities_by_key[key]
            current.evidence_refs = list(dict.fromkeys([*current.evidence_refs, *entity.evidence_refs]))[:40]
            current.aliases = list(dict.fromkeys([*current.aliases, *entity.aliases]))[:20]
            current.confidence = max(current.confidence, entity.confidence)
            current.attributes.update(entity.attributes)
        else:
            entities_by_key[key] = entity
    result.entities = list(entities_by_key.values())

    cleaned_claims_by_sig: dict[str, ClaimCandidate] = {}
    for raw in result.claims:
        claim = _clean_claim(raw)
        if claim is None:
            continue
        payload = claim.model_dump()
        payload.pop("evidence_refs", None)
        payload.pop("confidence", None)
        payload.pop("authority_score", None)
        sig = json.dumps(payload, ensure_ascii=False, sort_keys=True, default=str)
        if sig in cleaned_claims_by_sig:
            current = cleaned_claims_by_sig[sig]
            current.evidence_refs = list(dict.fromkeys([*current.evidence_refs, *claim.evidence_refs]))[:40]
            current.confidence = max(current.confidence, claim.confidence)
            if claim.authority_score is not None:
                current.authority_score = max(float(current.authority_score or 0.0), float(claim.authority_score))
        else:
            cleaned_claims_by_sig[sig] = claim
    result.claims = list(cleaned_claims_by_sig.values())

    cleaned_relations: list[RelationCandidate] = []
    seen_relations: set[str] = set()
    for raw in result.relations:
        relation = _clean_relation(raw)
        if relation is None:
            continue
        sig = f"{relation.source_key}|{relation.relation_type}|{relation.target_key}"
        if sig in seen_relations:
            continue
        seen_relations.add(sig)
        cleaned_relations.append(relation)
    result.relations = cleaned_relations

    result.unknowns = list(dict.fromkeys(str(v).strip() for v in result.unknowns if str(v).strip()))[:80]
    result.contradictions = list(dict.fromkeys(str(v).strip() for v in result.contradictions if str(v).strip()))[:80]
    result.warnings = list(dict.fromkeys(str(v).strip() for v in result.warnings if str(v).strip()))[:80]
    return result


def bench_role(source_role: str) -> str:
    return ROLE_ALIASES.get(source_role, source_role)


def result_to_bench_fragment(result: FileAnalysisResult, source_label: str) -> dict[str, Any]:
    """Converte uma análise de arquivo no contrato parcial do NAVE IQ Bench."""
    key_map = {entity.key: entity for entity in result.entities}
    entities = [
        {
            "id": entity.key,
            "type": entity.entity_type,
            "canonical_name": entity.canonical_name,
            "evidence_refs": [f"{source_label}:{ref}" for ref in entity.evidence_refs],
        }
        for entity in result.entities
    ]
    claims = []
    for claim in result.claims:
        row: dict[str, Any] = {
            "subject": claim.subject_key,
            "predicate": claim.predicate,
            "evidence_refs": [f"{source_label}:{ref}" for ref in claim.evidence_refs],
        }
        for field in ("value_text", "value_numeric", "value_boolean", "value_date", "value_timestamp", "value_json", "currency", "unit"):
            value = getattr(claim, field, None)
            if value is not None:
                row[field] = value
        if claim.object_key:
            row["object"] = claim.object_key
        claims.append(row)
    relations = [
        {
            "source": relation.source_key,
            "relation": relation.relation_type,
            "target": relation.target_key,
            "evidence_refs": [f"{source_label}:{ref}" for ref in relation.evidence_refs],
        }
        for relation in result.relations
    ]
    return {
        "source_role": bench_role(result.source_role),
        "entities": entities,
        "claims": claims,
        "relations": relations,
        "semantic_analysis_ran": result.semantic_analysis_ran,
        "warnings": list(result.warnings),
    }
