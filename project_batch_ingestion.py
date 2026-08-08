from __future__ import annotations

import hashlib
import io
import json
import mimetypes
import re
import unicodedata
from dataclasses import asdict, dataclass
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence

PROJECT_FILES_BUCKET = "nave-project-files"
WORKFLOW_VERSION = "28.1.7"
MAX_TEXT_CHARS = 60000
MAX_FILE_BYTES = 100 * 1024 * 1024

ROLE_LABELS: dict[str, str] = {
    "briefing_original": "Briefing original",
    "proposal_presentation": "Apresentação de proposta",
    "detailed_costs": "Planilha detalhada de custos",
    "preliminary_budget": "Estudo preliminar de verba",
    "final_presentation": "Apresentação final",
    "feedback_approval": "Feedback / aprovação",
    "post_event_report": "Relatório pós-evento / encerramento",
    "supplier_reference": "Fornecedor / referência técnica",
    "complementary_document": "Documento complementar",
}
LABEL_TO_ROLE = {label: key for key, label in ROLE_LABELS.items()}

TARGET_SECTIONS: dict[str, tuple[str, ...]] = {
    "briefing_original": (
        "Briefing original",
        "Diagnóstico e recomendações",
    ),
    "proposal_presentation": (
        "Estratégia e conceito",
        "Cenografia e ativações",
        "Brindes e press kits",
        "Fornecedores e referências",
        "Documentos",
    ),
    "detailed_costs": (
        "Orçamento e aderência",
        "Diagnóstico e recomendações",
        "Documentos",
    ),
    "preliminary_budget": (
        "Orçamento e aderência",
        "Diagnóstico e recomendações",
        "Documentos",
    ),
    "final_presentation": (
        "Estratégia e conceito",
        "Cenografia e ativações",
        "Brindes e press kits",
        "Fornecedores e referências",
        "Apresentações finais",
        "Documentos",
    ),
    "feedback_approval": (
        "Feedbacks e aprovações",
        "Diagnóstico e recomendações",
        "Documentos",
    ),
    "post_event_report": (
        "Resultados e aprendizados",
        "Diagnóstico e recomendações",
        "Documentos",
    ),
    "supplier_reference": (
        "Fornecedores e referências",
        "Documentos",
    ),
    "complementary_document": ("Documentos",),
}

ALLOWED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".ppt", ".xlsx", ".xlsm", ".xls",
    ".csv", ".txt", ".md", ".eml", ".msg",
}

MIME_BY_EXTENSION = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt": "application/vnd.ms-powerpoint",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xlsm": "application/vnd.ms-excel.sheet.macroEnabled.12",
    ".xls": "application/vnd.ms-excel",
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".eml": "message/rfc822",
    ".msg": "application/vnd.ms-outlook",
}

_ROLE_TERMS: dict[str, tuple[tuple[str, float], ...]] = {
    "briefing_original": (
        ("briefing", 7.0), ("debriefing", 6.0), ("brief interno", 6.0),
        ("objetivo e desafio", 3.0), ("publico alvo", 2.5),
        ("entregaveis", 2.2), ("obrigatoriedades", 2.5),
        ("mandatorios", 2.4), ("restricoes", 1.6),
    ),
    "proposal_presentation": (
        ("apresentacao", 2.5), ("proposta", 5.2), ("conceito criativo", 3.2),
        ("estrategia", 2.0), ("cenografia", 2.0), ("ativacoes", 2.0),
        ("experiencias", 1.2), ("defesa", 2.5), ("paper", 2.3),
    ),
    "detailed_costs": (
        ("planilha de custos", 8.0), ("custos", 4.5), ("custo", 2.0),
        ("orcamento detalhado", 6.0), ("quantidade", 1.0),
        ("valor unitario", 3.0), ("valor total", 2.5), ("fee", 2.0),
        ("impostos", 1.5), ("fornecedor", 0.8),
    ),
    "preliminary_budget": (
        ("estudo de verba", 8.0), ("verba preliminar", 8.0),
        ("estimativa preliminar", 7.0), ("budget preliminar", 7.0),
        ("pre orcamento", 6.0), ("pre-orcamento", 6.0),
        ("estimativa de custo", 5.0), ("verba", 2.0),
    ),
    "final_presentation": (
        ("apresentacao final", 8.0), ("proposta final", 7.0),
        ("final aprovado", 7.0), ("versao final", 6.0),
        ("v final", 4.0), ("final", 1.4), ("aprovada", 2.0),
    ),
    "feedback_approval": (
        ("feedback", 7.0), ("aprovacao", 6.0), ("aprovado", 5.0),
        ("retorno do cliente", 6.0), ("alteracoes solicitadas", 5.0),
        ("comentarios do cliente", 5.0), ("ajustes", 1.2),
    ),
    "post_event_report": (
        ("pos evento", 7.0), ("pos-evento", 7.0), ("post event", 7.0),
        ("relatorio final", 6.0), ("relatorio de encerramento", 8.0),
        ("report final", 7.0), ("resultados", 2.2), ("aprendizados", 3.0),
        ("kpis", 2.5), ("executado", 1.5), ("realizado", 1.0),
    ),
    "supplier_reference": (
        ("fornecedor", 4.0), ("cotacao", 4.0), ("proposta comercial", 5.0),
        ("ficha tecnica", 5.0), ("catalogo", 4.0), ("portfolio", 3.0),
        ("lead time", 2.0), ("prazo de entrega", 2.0), ("condicoes comerciais", 3.0),
    ),
}

_SPREADSHEET_EXTENSIONS = {".xlsx", ".xlsm", ".xls", ".csv"}


def _spreadsheet_cost_identity(
    *,
    extension: str,
    normalized_name: str,
    normalized_text: str,
) -> tuple[str | None, list[str]]:
    """Resolve o papel financeiro antes de termos incidentais do conteúdo.

    Planilhas humanas de orçamento costumam carregar textos como ``pós-evento`` ou
    ``relatório final`` dentro de uma linha de escopo. Esses termos descrevem um
    serviço, não a identidade documental da planilha. O formato + a estrutura da
    tabela devem pesar mais do que palavras encontradas em células isoladas.
    """
    if extension not in _SPREADSHEET_EXTENSIONS:
        return None, []

    haystack = f" {normalized_name} {normalized_text} "
    preliminary_signals = (
        "estudo de verba",
        "verba preliminar",
        "budget preliminar",
        "pre orcamento",
        "pre-orcamento",
        "estimativa preliminar",
    )
    if any(signal in normalized_name for signal in preliminary_signals):
        return "preliminary_budget", ["nome identifica estudo preliminar de verba"]

    detailed_groups = (
        ("tipo faturamento", "valor unit", "valor total"),
        ("total com honorarios", "honorarios", "encargos"),
        ("abertura de custos", "impostos", "total geral dos servicos"),
    )
    if any(all(signal in haystack for signal in group) for group in detailed_groups):
        return "detailed_costs", ["estrutura tabular de abertura de custos"]

    # Resumo de categorias com valores, honorários e impostos também é uma
    # evidência forte de orçamento detalhado mesmo sem cabeçalho convencional.
    financial_signals = sum(
        signal in haystack
        for signal in (
            "planejamento e producao",
            "infraestrutura",
            "comunicacao",
            "brindes",
            "artistico",
            "staff",
            "honorarios",
            "impostos",
            "total geral",
        )
    )
    if financial_signals >= 6 and ("honorarios" in haystack or "valor total" in haystack):
        return "detailed_costs", ["categorias financeiras e totais de produção"]

    if "budget" in haystack and "estudo de verba" in haystack:
        return "preliminary_budget", ["estrutura de budget preliminar"]
    return None, []


_NAME_STOPWORDS = {
    "briefing", "brief", "debriefing", "interno", "interna", "voe", "ideias",
    "apresentacao", "apresentação", "proposta", "final", "compressed", "compactado",
    "custos", "custo", "orcamento", "orçamento", "budget", "verba", "planilha",
    "relatorio", "relatório", "report", "feedback", "aprovacao", "aprovação",
    "versao", "versão", "rev", "revisao", "revisão", "arquivo", "documento",
    "estudo", "preliminar", "v1", "v2", "v3", "v4", "v5", "v6", "v7", "v8", "v9",
}


@dataclass
class PreparedDocument:
    name: str
    extension: str
    mime_type: str
    data: bytes
    sha256: str
    file_size_bytes: int
    text_excerpt: str
    page_count: int | None
    role: str
    role_confidence: float
    role_reasons: list[str]

    @property
    def role_label(self) -> str:
        return ROLE_LABELS.get(self.role, ROLE_LABELS["complementary_document"])

    @property
    def target_sections(self) -> tuple[str, ...]:
        return TARGET_SECTIONS.get(self.role, TARGET_SECTIONS["complementary_document"])

    def metadata_for_json(self) -> dict[str, Any]:
        data = asdict(self)
        data.pop("data", None)
        data["role_label"] = self.role_label
        data["target_sections"] = list(self.target_sections)
        return data


@dataclass
class ProjectCandidate:
    project_id: str
    project_name: str
    client_brand: str | None
    event_name: str | None
    score: float
    confidence: str
    reasons: list[str]
    conflicts: list[str]


class ProjectBatchError(RuntimeError):
    pass


def normalize_text(value: Any) -> str:
    text = unicodedata.normalize("NFKD", str(value or ""))
    text = "".join(ch for ch in text if not unicodedata.combining(ch))
    text = text.casefold()
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def safe_filename(value: str) -> str:
    path = Path(str(value or "arquivo"))
    stem = normalize_text(path.stem).replace(" ", "_")[:110] or "arquivo"
    suffix = path.suffix.lower()
    return f"{stem}{suffix}"


def infer_mime_type(name: str, supplied: str | None = None) -> str:
    if supplied and supplied != "application/octet-stream":
        return str(supplied)
    suffix = Path(name).suffix.lower()
    return MIME_BY_EXTENSION.get(suffix) or mimetypes.guess_type(name)[0] or "application/octet-stream"


def _clip(value: Any, size: int = MAX_TEXT_CHARS) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text[:size]


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return data.decode(encoding, errors="strict")
        except Exception:
            continue
    return data.decode("utf-8", errors="ignore")


def _pdf_text(data: bytes) -> tuple[str, int | None]:
    try:
        import pymupdf  # type: ignore
        pdf = pymupdf.open(stream=data, filetype="pdf")
    except Exception:
        try:
            import fitz  # type: ignore
            pdf = fitz.open(stream=data, filetype="pdf")
        except Exception:
            return "", None
    try:
        chunks: list[str] = []
        for index in range(min(len(pdf), 30)):
            try:
                chunks.append(pdf[index].get_text("text"))
            except Exception:
                continue
            if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                break
        return _clip("\n".join(chunks)), len(pdf)
    finally:
        pdf.close()


def _docx_text(data: bytes) -> str:
    try:
        from docx import Document  # type: ignore
        document = Document(io.BytesIO(data))
        chunks = [p.text for p in document.paragraphs if p.text]
        for table in document.tables[:20]:
            for row in table.rows[:60]:
                chunks.append(" | ".join(cell.text for cell in row.cells if cell.text))
                if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                    break
        return _clip("\n".join(chunks))
    except Exception:
        return ""


def _pptx_text(data: bytes) -> tuple[str, int | None]:
    try:
        from pptx import Presentation  # type: ignore
        presentation = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for slide in presentation.slides[:80]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and getattr(shape, "text", None):
                    chunks.append(str(shape.text))
                if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                    break
            if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                break
        return _clip("\n".join(chunks)), len(presentation.slides)
    except Exception:
        return "", None


def _xlsx_text(data: bytes, extension: str) -> str:
    if extension in {".xlsx", ".xlsm"}:
        try:
            from openpyxl import load_workbook  # type: ignore
            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            chunks: list[str] = []
            for sheet in workbook.worksheets[:8]:
                chunks.append(f"Planilha: {sheet.title}")
                for row in sheet.iter_rows(min_row=1, max_row=90, max_col=24, values_only=True):
                    values = [str(value) for value in row if value not in (None, "")]
                    if values:
                        chunks.append(" | ".join(values))
                    if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                        break
                if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
                    break
            workbook.close()
            return _clip("\n".join(chunks))
        except Exception:
            return ""
    try:
        import pandas as pd  # type: ignore
        sheets = pd.read_excel(io.BytesIO(data), sheet_name=None, nrows=80)
        chunks: list[str] = []
        for sheet_name, frame in list(sheets.items())[:8]:
            chunks.append(f"Planilha: {sheet_name}")
            chunks.append(frame.head(80).astype(str).to_csv(index=False))
        return _clip("\n".join(chunks))
    except Exception:
        return ""


def _eml_text(data: bytes) -> str:
    try:
        from email import policy
        from email.parser import BytesParser
        message = BytesParser(policy=policy.default).parsebytes(data)
        chunks = [str(message.get("subject") or "")]
        if message.is_multipart():
            for part in message.walk():
                if part.get_content_type() == "text/plain":
                    chunks.append(str(part.get_content()))
        else:
            chunks.append(str(message.get_content()))
        return _clip("\n".join(chunks))
    except Exception:
        return ""


def extract_text(data: bytes, name: str) -> tuple[str, int | None]:
    extension = Path(name).suffix.lower()
    if extension == ".pdf":
        return _pdf_text(data)
    if extension == ".docx":
        return _docx_text(data), None
    if extension == ".pptx":
        return _pptx_text(data)
    if extension in {".xlsx", ".xlsm", ".xls"}:
        return _xlsx_text(data, extension), None
    if extension == ".eml":
        return _eml_text(data), None
    if extension in {".csv", ".txt", ".md"}:
        return _clip(_decode_text(data)), None
    return "", None


def _term_score(haystack: str, role: str) -> tuple[float, list[str]]:
    score = 0.0
    reasons: list[str] = []
    for term, weight in _ROLE_TERMS.get(role, ()):
        if normalize_text(term) in haystack:
            score += weight
            if len(reasons) < 4:
                reasons.append(term)
    return score, reasons


def classify_document(name: str, text_excerpt: str) -> tuple[str, float, list[str]]:
    normalized_name = normalize_text(Path(name).stem)
    normalized_text = normalize_text(text_excerpt[:MAX_TEXT_CHARS])
    name_haystack = f" {normalized_name} "
    full_haystack = f" {normalized_name} {normalized_text} "
    extension = Path(name).suffix.lower()

    spreadsheet_role, spreadsheet_reasons = _spreadsheet_cost_identity(
        extension=extension,
        normalized_name=normalized_name,
        normalized_text=normalized_text,
    )
    if spreadsheet_role:
        return spreadsheet_role, 0.98, spreadsheet_reasons

    scores: dict[str, float] = {role: 0.0 for role in ROLE_LABELS}
    reasons: dict[str, list[str]] = {role: [] for role in ROLE_LABELS}

    for role in _ROLE_TERMS:
        content_score, content_reasons = _term_score(full_haystack, role)
        name_score, name_reasons = _term_score(name_haystack, role)
        scores[role] += content_score + name_score * 0.7
        reasons[role].extend([f"termo: {term}" for term in (name_reasons + content_reasons)[:5]])

    if extension in {".xlsx", ".xlsm", ".xls", ".csv"}:
        scores["detailed_costs"] += 4.0
        reasons["detailed_costs"].append("formato de planilha")
    elif extension in {".pptx", ".ppt"}:
        scores["proposal_presentation"] += 2.8
        reasons["proposal_presentation"].append("formato de apresentação")
    elif extension in {".docx", ".txt", ".md"}:
        scores["briefing_original"] += 1.4
        reasons["briefing_original"].append("formato textual")
    elif extension == ".eml":
        scores["feedback_approval"] += 2.2
        reasons["feedback_approval"].append("mensagem de e-mail")

    # Regras de precedência para evitar que todo arquivo final seja tratado só como proposta.
    specific_roles = (
        "post_event_report", "feedback_approval", "preliminary_budget",
        "detailed_costs", "briefing_original", "final_presentation",
        "supplier_reference", "proposal_presentation",
    )
    ranked = sorted(specific_roles, key=lambda role: scores[role], reverse=True)
    winner = ranked[0]
    top = scores[winner]
    second = scores[ranked[1]] if len(ranked) > 1 else 0.0

    if top < 2.2:
        return "complementary_document", 0.45, ["sem sinal forte; revisar manualmente"]

    if winner == "proposal_presentation" and scores["final_presentation"] >= max(5.0, top * 0.72):
        winner = "final_presentation"
        top = scores[winner]
        second = max(score for role, score in scores.items() if role != winner)

    margin = max(0.0, top - second)
    confidence = min(0.98, 0.50 + min(top, 12.0) / 30.0 + min(margin, 8.0) / 24.0)
    if margin < 1.2:
        confidence = min(confidence, 0.68)
        reasons[winner].append("classificação próxima de outra categoria")
    return winner, round(confidence, 3), list(dict.fromkeys(reasons[winner]))[:6]


def prepare_document(name: str, data: bytes, mime_type: str | None = None) -> PreparedDocument:
    extension = Path(name).suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise ProjectBatchError(f"Formato não suportado nesta etapa: {name}")
    if len(data) > MAX_FILE_BYTES:
        raise ProjectBatchError(f"{name} ultrapassa o limite de 100 MB da NAVE.")
    excerpt, page_count = extract_text(data, name)
    role, confidence, reasons = classify_document(name, excerpt)
    return PreparedDocument(
        name=name,
        extension=extension,
        mime_type=infer_mime_type(name, mime_type),
        data=data,
        sha256=hashlib.sha256(data).hexdigest(),
        file_size_bytes=len(data),
        text_excerpt=excerpt[:12000],
        page_count=page_count,
        role=role,
        role_confidence=confidence,
        role_reasons=reasons,
    )


def prepare_documents(files: Iterable[tuple[str, bytes, str | None]]) -> list[PreparedDocument]:
    result: list[PreparedDocument] = []
    for name, data, mime_type in files:
        result.append(prepare_document(name, data, mime_type))
    return result


def infer_project_name(documents: Sequence[PreparedDocument]) -> str:
    if not documents:
        return ""

    token_sets: list[list[str]] = []
    for document in documents:
        stem = re.sub(r"[_\-–—]+", " ", Path(document.name).stem)
        tokens = [token for token in re.split(r"\s+", stem.strip()) if token]
        cleaned = [
            token for token in tokens
            if normalize_text(token) not in _NAME_STOPWORDS
            and not re.fullmatch(r"v\d+(?:\.\d+)*", normalize_text(token))
        ]
        token_sets.append(cleaned)

    if len(token_sets) == 1:
        return " ".join(token_sets[0][:10]).strip()

    threshold = max(2, (len(token_sets) + 1) // 2)
    counts: dict[str, int] = {}
    display: dict[str, str] = {}
    for tokens in token_sets:
        seen: set[str] = set()
        for token in tokens:
            key = normalize_text(token)
            if not key or key in seen:
                continue
            seen.add(key)
            counts[key] = counts.get(key, 0) + 1
            display.setdefault(key, token)

    common = [key for key, count in counts.items() if count >= threshold]
    if common:
        ordered: list[str] = []
        first = token_sets[0]
        for token in first:
            key = normalize_text(token)
            if key in common and key not in {normalize_text(x) for x in ordered}:
                ordered.append(display[key])
        if ordered:
            return " ".join(ordered[:10]).strip()

    return " ".join(token_sets[0][:8]).strip()


def _extract_years(text: str) -> set[int]:
    return {int(item) for item in re.findall(r"\b(20\d{2})\b", str(text or ""))}


def _extract_named_editions(text: str) -> set[str]:
    normalized = normalize_text(text)
    found = set()
    for pattern in (
        r"\bplaneja\s*(\d{2,4})\b",
        r"\bccxp\s*(\d{2,4})\b",
        r"\bedicao\s*(\d{1,4})\b",
    ):
        found.update(re.findall(pattern, normalized))
    return found


def fetch_projects(client: Any) -> list[dict[str, Any]]:
    response = (
        client.table("projects")
        .select("id,project_name,client_brand,event_name,status,event_date,location_city,location_state,raw_data,updated_at")
        .order("updated_at", desc=True)
        .limit(2500)
        .execute()
    )
    return [dict(row) for row in (getattr(response, "data", None) or [])]


def rank_project_candidates(
    projects: Sequence[Mapping[str, Any]],
    *,
    project_name: str,
    client_brand: str | None = None,
    event_name: str | None = None,
    limit: int = 5,
) -> list[ProjectCandidate]:
    query_name = normalize_text(project_name)
    query_client = normalize_text(client_brand)
    query_event = normalize_text(event_name)
    query_years = _extract_years(project_name + " " + str(event_name or ""))
    query_editions = _extract_named_editions(project_name + " " + str(event_name or ""))
    candidates: list[ProjectCandidate] = []

    if not query_name and not query_client and not query_event:
        return []

    for row in projects:
        candidate_name = normalize_text(row.get("project_name"))
        candidate_client = normalize_text(row.get("client_brand"))
        candidate_event = normalize_text(row.get("event_name"))
        raw_data = row.get("raw_data") if isinstance(row.get("raw_data"), dict) else {}
        combined = " ".join(
            str(value or "")
            for value in (
                row.get("project_name"), row.get("event_name"), row.get("event_date"),
                raw_data.get("reference_year"), raw_data.get("edition"), raw_data.get("year"),
            )
        )
        candidate_years = _extract_years(combined)
        candidate_editions = _extract_named_editions(combined)
        reasons: list[str] = []
        conflicts: list[str] = []

        name_similarity = SequenceMatcher(None, query_name, candidate_name).ratio() if query_name and candidate_name else 0.0
        score = name_similarity * 0.72
        if name_similarity >= 0.82:
            reasons.append("nome do projeto muito semelhante")
        elif name_similarity >= 0.62:
            reasons.append("nome do projeto parcialmente semelhante")

        if query_client and candidate_client:
            client_similarity = SequenceMatcher(None, query_client, candidate_client).ratio()
            score += client_similarity * 0.18
            if client_similarity >= 0.88:
                reasons.append("cliente / marca coincide")
            elif client_similarity < 0.45:
                score -= 0.22
                conflicts.append("cliente / marca diverge")

        if query_event and candidate_event:
            event_similarity = SequenceMatcher(None, query_event, candidate_event).ratio()
            score += event_similarity * 0.10
            if event_similarity >= 0.80:
                reasons.append("evento coincide")

        if query_years and candidate_years and query_years.isdisjoint(candidate_years):
            score -= 0.35
            conflicts.append("ano de referência diverge")
        if query_editions and candidate_editions and query_editions.isdisjoint(candidate_editions):
            score -= 0.45
            conflicts.append("edição diverge")

        score = max(0.0, min(1.0, score))
        if score >= 0.78 and not conflicts:
            confidence = "alta"
        elif score >= 0.56:
            confidence = "média"
        else:
            confidence = "baixa"

        if score >= 0.28:
            candidates.append(ProjectCandidate(
                project_id=str(row.get("id")),
                project_name=str(row.get("project_name") or "Projeto sem nome"),
                client_brand=str(row.get("client_brand") or "").strip() or None,
                event_name=str(row.get("event_name") or "").strip() or None,
                score=round(score, 3),
                confidence=confidence,
                reasons=reasons,
                conflicts=conflicts,
            ))

    candidates.sort(key=lambda item: item.score, reverse=True)
    return candidates[:limit]


def _ensure_bucket(client: Any) -> None:
    try:
        buckets = client.storage.list_buckets()
        names = {
            str(getattr(bucket, "name", "") or (bucket.get("name") if isinstance(bucket, dict) else ""))
            for bucket in (buckets or [])
        }
        if PROJECT_FILES_BUCKET in names:
            return
    except Exception:
        # O bucket já é parte do workspace V27; caso a listagem não esteja disponível,
        # tentamos criar e tratamos "already exists" como sucesso.
        pass

    try:
        client.storage.create_bucket(PROJECT_FILES_BUCKET, options={"public": False})
    except Exception as exc:
        message = str(exc).casefold()
        if "already" not in message and "exist" not in message and "duplicate" not in message:
            # Ainda pode existir sem permissão de listagem. A tentativa real de upload
            # produzirá a mensagem definitiva se o bucket estiver indisponível.
            return


def _upload_bytes(client: Any, *, path: str, data: bytes, mime_type: str) -> None:
    try:
        client.storage.from_(PROJECT_FILES_BUCKET).upload(
            path,
            data,
            file_options={"content-type": mime_type, "upsert": "false"},
        )
    except TypeError:
        client.storage.from_(PROJECT_FILES_BUCKET).upload(
            path,
            data,
            {"content-type": mime_type, "upsert": "false"},
        )


def _project_raw_data(client: Any, project_id: str) -> dict[str, Any]:
    try:
        response = client.table("projects").select("raw_data").eq("id", project_id).limit(1).execute()
        rows = getattr(response, "data", None) or []
        if rows and isinstance(rows[0].get("raw_data"), dict):
            return dict(rows[0]["raw_data"])
    except Exception:
        pass
    return {}


def _create_project(
    client: Any,
    *,
    project_name: str,
    client_brand: str | None,
    event_name: str | None,
) -> str:
    clean_name = str(project_name or "").strip()
    if not clean_name:
        raise ProjectBatchError("Informe o nome do projeto antes de importar.")
    payload = {
        "project_name": clean_name,
        "normalized_name": normalize_text(clean_name),
        "client_brand": str(client_brand or "").strip() or None,
        "event_name": str(event_name or "").strip() or None,
        "status": "rascunho",
        "raw_data": {
            "source": "v28_1_multi_document_ingestion",
            "created_from_project_bundle": True,
            "workflow_version": WORKFLOW_VERSION,
        },
    }
    response = client.table("projects").insert(payload).execute()
    rows = getattr(response, "data", None) or []
    if not rows:
        raise ProjectBatchError("O Supabase não confirmou a criação do projeto.")
    return str(rows[0]["id"])


def _existing_by_sha(client: Any, project_id: str, sha256: str) -> dict[str, Any] | None:
    try:
        response = (
            client.table("source_files")
            .select("id,storage_bucket,storage_path,file_name,sha256")
            .eq("project_id", project_id)
            .eq("sha256", sha256)
            .limit(1)
            .execute()
        )
        rows = getattr(response, "data", None) or []
        return dict(rows[0]) if rows else None
    except Exception:
        return None


def _source_file_payload(
    *,
    import_id: str,
    project_id: str,
    document: PreparedDocument,
    role: str,
    confidence: float,
    reasons: Sequence[str],
    storage_path: str | None,
    batch_order: int,
    duplicate: bool,
) -> dict[str, Any]:
    targets = list(TARGET_SECTIONS.get(role, TARGET_SECTIONS["complementary_document"]))
    return {
        "import_id": import_id,
        "project_id": project_id,
        "file_name": document.name,
        "mime_type": document.mime_type,
        "storage_bucket": PROJECT_FILES_BUCKET if storage_path else None,
        "storage_path": storage_path,
        "page_count": document.page_count,
        "sha256": document.sha256,
        "document_role": role,
        "role_confidence": max(0.0, min(1.0, float(confidence))),
        "role_reasons": list(reasons)[:12],
        "target_sections": targets,
        "text_excerpt": document.text_excerpt[:12000] or None,
        "file_size_bytes": document.file_size_bytes,
        "batch_order": batch_order,
        "processing_status": "skipped_duplicate" if duplicate else "ready",
        "metadata": {
            "workflow_version": WORKFLOW_VERSION,
            "role_label": ROLE_LABELS.get(role),
            "duplicate_in_project": duplicate,
        },
    }


def save_project_bundle(
    client: Any,
    *,
    documents: Sequence[PreparedDocument],
    role_overrides: Mapping[str, str] | None,
    include_sha256: set[str] | None,
    project_name: str,
    client_brand: str | None,
    event_name: str | None,
    existing_project_id: str | None = None,
    match_context: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    selected = [
        document for document in documents
        if include_sha256 is None or document.sha256 in include_sha256
    ]
    if not selected:
        raise ProjectBatchError("Selecione pelo menos um arquivo para importar.")

    _ensure_bucket(client)
    created_project = False
    project_id = str(existing_project_id or "").strip()
    if not project_id:
        project_id = _create_project(
            client,
            project_name=project_name,
            client_brand=client_brand,
            event_name=event_name,
        )
        created_project = True

    roles_for_document: dict[str, str] = {}
    for document in selected:
        role = str((role_overrides or {}).get(document.sha256) or document.role)
        if role not in ROLE_LABELS:
            role = "complementary_document"
        roles_for_document[document.sha256] = role

    classification = {
        "workflow": "project_bundle",
        "workflow_version": WORKFLOW_VERSION,
        "project_name": str(project_name or "").strip() or None,
        "client_brand": str(client_brand or "").strip() or None,
        "event_name": str(event_name or "").strip() or None,
        "existing_project": bool(existing_project_id),
        "match_context": dict(match_context or {}),
        "documents": [
            {
                **document.metadata_for_json(),
                "role": roles_for_document[document.sha256],
                "role_label": ROLE_LABELS[roles_for_document[document.sha256]],
                "target_sections": list(TARGET_SECTIONS[roles_for_document[document.sha256]]),
            }
            for document in selected
        ],
    }
    source_files_summary = [
        {
            "file_name": document.name,
            "mime_type": document.mime_type,
            "sha256": document.sha256,
            "file_size_bytes": document.file_size_bytes,
            "document_role": roles_for_document[document.sha256],
        }
        for document in selected
    ]
    import_payload = {
        "document_type": "project_bundle",
        "destination_base": "projects",
        "document_title": str(project_name or "").strip() or "Conjunto de documentos de projeto",
        "project_id": project_id,
        "source_files": source_files_summary,
        "classification": classification,
        "original_payload": {"workflow_version": WORKFLOW_VERSION},
        "warnings": [],
        "status": "processando",
        "imported_records": 0,
    }

    uploaded_paths: list[str] = []
    import_id: str | None = None
    saved_rows: list[dict[str, Any]] = []
    duplicate_count = 0
    try:
        inserted_import = client.table("imports").insert(import_payload).execute()
        import_rows = getattr(inserted_import, "data", None) or []
        if not import_rows:
            raise ProjectBatchError("Não foi possível abrir o lote de importação.")
        import_id = str(import_rows[0]["id"])

        seen_batch_hashes: dict[str, str] = {}
        for order, document in enumerate(selected, start=1):
            existing = _existing_by_sha(client, project_id, document.sha256)
            existing_path = str((existing or {}).get("storage_path") or "").strip()
            batch_path = seen_batch_hashes.get(document.sha256)
            duplicate = bool(existing_path or batch_path)
            if duplicate:
                duplicate_count += 1
                storage_path = existing_path or batch_path
            else:
                storage_path = (
                    f"projects/{project_id}/imports/{import_id}/"
                    f"{order:02d}_{document.sha256[:12]}_{safe_filename(document.name)}"
                )
                _upload_bytes(
                    client,
                    path=storage_path,
                    data=document.data,
                    mime_type=document.mime_type,
                )
                uploaded_paths.append(storage_path)
                seen_batch_hashes[document.sha256] = storage_path

            payload = _source_file_payload(
                import_id=import_id,
                project_id=project_id,
                document=document,
                role=roles_for_document[document.sha256],
                confidence=document.role_confidence,
                reasons=document.role_reasons,
                storage_path=storage_path,
                batch_order=order,
                duplicate=duplicate,
            )
            inserted_file = client.table("source_files").insert(payload).execute()
            rows = getattr(inserted_file, "data", None) or []
            if rows:
                saved_rows.append(dict(rows[0]))

        # V28.1.1: o lote deixa de apenas classificar/rotear documentos e
        # passa a materializá-los nas estruturas que o workspace realmente lê.
        # A materialização é idempotente por projeto + hash e, portanto, também
        # pode ser executada sobre arquivos reaproveitados sem criar duplicatas.
        from project_bundle_materializer import materialize_new_source_files

        documents_by_sha = {document.sha256: document for document in selected}
        materialization_results = materialize_new_source_files(
            client,
            saved_rows,
            documents_by_sha=documents_by_sha,
        )
        workspace_materialized = sum(
            1 for item in materialization_results
            if str(item.get("status") or "") != "error"
        )
        workspace_errors = sum(
            1 for item in materialization_results
            if str(item.get("status") or "") == "error"
        )
        workspace_warnings = [
            warning
            for item in materialization_results
            for warning in (item.get("warnings") or [])
            if str(warning).strip()
        ]

        import_warnings = []
        if duplicate_count:
            import_warnings.append(
                f"{duplicate_count} arquivo(s) já existiam no projeto e foram reaproveitados."
            )
        if workspace_errors:
            import_warnings.append(
                f"{workspace_errors} arquivo(s) tiveram erro ao incorporar conteúdo ao workspace."
            )
        if workspace_warnings:
            import_warnings.extend(str(item)[:900] for item in workspace_warnings[:20])

        try:
            client.table("imports").update({
                "status": "projeto_importado" if not workspace_errors else "projeto_importado_com_alertas",
                "imported_records": len(saved_rows),
                "warnings": import_warnings,
            }).eq("id", import_id).execute()
        except Exception as exc:
            workspace_warnings.append(
                f"O lote foi incorporado ao workspace, mas o resumo técnico do import não pôde ser atualizado: {exc}"
            )

        try:
            current_raw = _project_raw_data(client, project_id)
            current_raw.update({
                "last_project_bundle_import_id": import_id,
                "last_project_bundle_at_version": WORKFLOW_VERSION,
                "last_project_bundle_documents": len(selected),
                "last_project_bundle_roles": sorted(set(roles_for_document.values())),
                "last_project_bundle_materialized": workspace_materialized,
                "last_project_bundle_errors": workspace_errors,
            })
            client.table("projects").update({"raw_data": current_raw}).eq("id", project_id).execute()
        except Exception as exc:
            workspace_warnings.append(
                f"O conteúdo foi incorporado, mas o resumo técnico do projeto não pôde ser atualizado: {exc}"
            )

        return {
            "status": "saved" if not workspace_errors else "saved_with_warnings",
            "project_id": project_id,
            "import_id": import_id,
            "created_project": created_project,
            "documents_saved": len(saved_rows),
            "duplicates_reused": duplicate_count,
            "roles": roles_for_document,
            "workspace_materialized": workspace_materialized,
            "workspace_errors": workspace_errors,
            "workspace_warnings": workspace_warnings[:40],
            "materialization_results": materialization_results,
        }

    except Exception as exc:
        if uploaded_paths:
            try:
                client.storage.from_(PROJECT_FILES_BUCKET).remove(uploaded_paths)
            except Exception:
                pass
        if import_id:
            try:
                client.table("imports").delete().eq("id", import_id).execute()
            except Exception:
                pass
        if created_project:
            try:
                client.table("projects").delete().eq("id", project_id).execute()
            except Exception:
                pass
        if isinstance(exc, ProjectBatchError):
            raise
        raise ProjectBatchError(
            "A importação foi interrompida e as alterações parciais foram revertidas. "
            f"Detalhe técnico: {exc}"
        ) from exc


def roles_summary(documents: Sequence[PreparedDocument]) -> dict[str, int]:
    result: dict[str, int] = {}
    for document in documents:
        result[document.role] = result.get(document.role, 0) + 1
    return result


def json_safe_prepared(documents: Sequence[PreparedDocument]) -> str:
    return json.dumps([item.metadata_for_json() for item in documents], ensure_ascii=False, indent=2)
